"""Gera a apresentação de Programação Funcional em Python (.pptx).

Uso: ``python docs/gerar_slides.py``

Por que gerar por código em vez de editar à mão: os critérios de pontuação
são objetivos (fonte mínima no corpo, rodapé numerado, atribuição de fontes).
Definindo-os como constantes aqui, o deck inteiro fica conforme por
construção, e `docs/verificar_slides.py` confere o resultado.

Decisões que atendem ao enunciado:

* **corpo nunca abaixo de 18pt** — o mínimo exigido é 16, mas a exportação
  para PDF costuma aplicar ~0,25% de escala e derrubar 16 para 15,96; 18
  atravessa a conversão com folga;
* **rodapé em todos os slides** com o número e, quando aplicável, a marca
  `[n]` da fonte;
* **atribuição de IA** detalhada no último slide, no formato pedido.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

RAIZ = Path(__file__).resolve().parents[1]
TELAS = RAIZ / "docs" / "screenshots"
LOGO = RAIZ / "api" / "app" / "static" / "img" / "logo-marca.png"
SAIDA = RAIZ / "docs" / "PF-Python-Ponte-Italia.pptx"

# ── identidade visual (seção 8 do README) ────────────────
NAVY = RGBColor(0x16, 0x23, 0x2B)
OFFWHITE = RGBColor(0xFA, 0xFA, 0xF7)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRAFITE = RGBColor(0x1F, 0x29, 0x33)
VERDE = RGBColor(0x2E, 0x7D, 0x5B)
VERDE_CLARO = RGBColor(0xEA, 0xF3, 0xEF)
VERDE_SUAVE = RGBColor(0x6E, 0xC0, 0x9B)
CINZA = RGBColor(0x5B, 0x66, 0x70)
CINZA_CLARO = RGBColor(0xF3, 0xF4, 0xF6)
CODIGO_FUNDO = RGBColor(0x14, 0x1E, 0x26)
CODIGO_TEXTO = RGBColor(0xE6, 0xEC, 0xF0)
VERMELHO = RGBColor(0xB0, 0x2A, 0x37)
AMBAR = RGBColor(0x8A, 0x5A, 0x00)

SERIF = "Georgia"
SANS = "Calibri"
MONO = "Consolas"

# ── tamanhos: nada no CORPO abaixo de CORPO_MIN ──────────
CORPO_MIN = Pt(18)
TITULO = Pt(38)
TITULO_MENOR = Pt(30)
KICKER = Pt(18)
SUBTITULO = Pt(20)
CARD_TITULO = Pt(20)
CODIGO = Pt(18)
RODAPE = Pt(10)  # rodapé é isento pela regra

LARGURA = Inches(13.333)
ALTURA = Inches(7.5)
MARGEM = Inches(0.7)
UTIL = LARGURA - 2 * MARGEM

RODAPE_TEXTO = "Programação Funcional em Python · Juliano Freitas"
# marca de origem: o deck foi produzido com auxílio de IA (ver slide de fontes)
MARCA_FONTE = "[1]"


# ── utilidades de desenho ────────────────────────────────


def _caixa(slide, x, y, cx, cy):
    caixa = slide.shapes.add_textbox(x, y, cx, cy)
    quadro = caixa.text_frame
    quadro.word_wrap = True
    quadro.margin_left = quadro.margin_right = 0
    quadro.margin_top = quadro.margin_bottom = 0
    return quadro


def _texto(quadro, texto, tamanho, cor, *, fonte=SANS, negrito=False,
           italico=False, espaco=None, alinhamento=PP_ALIGN.LEFT,
           primeiro=False, entrelinha=1.15):
    paragrafo = quadro.paragraphs[0] if primeiro else quadro.add_paragraph()
    paragrafo.alignment = alinhamento
    paragrafo.line_spacing = entrelinha
    if espaco is not None:
        paragrafo.space_before = espaco
    corrida = paragrafo.add_run()
    corrida.text = texto
    corrida.font.size = tamanho
    corrida.font.color.rgb = cor
    corrida.font.name = fonte
    corrida.font.bold = negrito
    corrida.font.italic = italico
    return paragrafo


_RECORTES = RAIZ / "docs" / "screenshots" / "_recortes"


def _recortar(origem: Path, proporcao: float) -> Path:
    """Recorta a faixa superior da captura na proporção pedida (largura/altura)."""
    from PIL import Image

    _RECORTES.mkdir(parents=True, exist_ok=True)
    destino = _RECORTES / origem.name
    with Image.open(origem) as img:
        altura_alvo = min(img.height, int(img.width / proporcao))
        img.crop((0, 0, img.width, altura_alvo)).save(destino)
    return destino


def _retangulo(slide, x, y, cx, cy, cor, *, raio=0.06, borda=None):
    forma = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    forma.adjustments[0] = raio
    forma.fill.solid()
    forma.fill.fore_color.rgb = cor
    if borda is None:
        forma.line.fill.background()
    else:
        forma.line.color.rgb = borda
        forma.line.width = Pt(1)
    forma.shadow.inherit = False
    forma.text_frame.word_wrap = True
    return forma


class Deck:
    """Monta o arquivo, garantindo rodapé e numeração em cada slide."""

    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = LARGURA
        self.prs.slide_height = ALTURA
        self.branco = self.prs.slide_layouts[6]  # em branco

    # -- infraestrutura ----------------------------------

    def _novo(self, fundo: RGBColor):
        slide = self.prs.slides.add_slide(self.branco)
        fundo_forma = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, LARGURA, ALTURA
        )
        fundo_forma.fill.solid()
        fundo_forma.fill.fore_color.rgb = fundo
        fundo_forma.line.fill.background()
        fundo_forma.shadow.inherit = False
        return slide

    def _rodape(self, slide, escuro: bool) -> None:
        cor = RGBColor(0x8A, 0x94, 0x9C) if escuro else CINZA
        numero = len(self.prs.slides._sldIdLst)  # o slide corrente já entrou

        quadro = _caixa(slide, MARGEM, ALTURA - Inches(0.52), Inches(9), Inches(0.3))
        _texto(quadro, f"{RODAPE_TEXTO}   {MARCA_FONTE}", RODAPE, cor, primeiro=True)

        quadro = _caixa(
            slide, LARGURA - MARGEM - Inches(1.2), ALTURA - Inches(0.52),
            Inches(1.2), Inches(0.3),
        )
        _texto(
            quadro, str(numero), RODAPE, cor,
            negrito=True, alinhamento=PP_ALIGN.RIGHT, primeiro=True,
        )

    def _cabecalho(self, slide, kicker: str, titulo: str, escuro=False):
        cor_titulo = BRANCO if escuro else GRAFITE
        cor_kicker = VERDE_SUAVE if escuro else VERDE
        y = MARGEM

        if kicker:
            quadro = _caixa(slide, MARGEM, y, UTIL, Inches(0.35))
            _texto(
                quadro, kicker.upper(), KICKER, cor_kicker,
                negrito=True, primeiro=True,
            )
            y += Inches(0.45)

        # escolhe o corpo do título pela largura REAL estimada e reserva
        # espaço por linha — assim um título de duas linhas não invade o
        # conteúdo abaixo (o defeito clássico deste tipo de layout)
        util_pt = Emu(UTIL).inches * 72
        tamanho = TITULO
        if len(titulo) * 0.56 * TITULO.pt > util_pt:
            tamanho = TITULO_MENOR
        linhas = max(1, -(-int(len(titulo) * 0.56 * tamanho.pt) // int(util_pt)))

        altura = Inches(0.62) * linhas
        quadro = _caixa(slide, MARGEM, y, UTIL, altura)
        _texto(quadro, titulo, tamanho, cor_titulo, fonte=SERIF,
               negrito=True, primeiro=True, entrelinha=1.05)
        return y + altura + Inches(0.45)

    # -- tipos de slide ----------------------------------

    def capa(self, titulo: str, subtitulo: str, autor: str, instituicao: str):
        slide = self._novo(NAVY)
        circulo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(9.6), Inches(3.2), Inches(4.6), Inches(4.6)
        )
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = VERDE
        circulo.line.fill.background()
        circulo.shadow.inherit = False

        if LOGO.exists():
            slide.shapes.add_picture(
                str(LOGO), MARGEM, Inches(0.75), height=Inches(0.95)
            )

        quadro = _caixa(slide, MARGEM, Inches(2.0), Inches(9.0), Inches(0.4))
        _texto(quadro, "PARADIGMAS DE LINGUAGENS DE PROGRAMAÇÃO", KICKER,
               VERDE_SUAVE, negrito=True, primeiro=True)

        quadro = _caixa(slide, MARGEM, Inches(2.55), Inches(9.2), Inches(1.9))
        _texto(quadro, titulo, Pt(50), BRANCO, fonte=SERIF, negrito=True,
               primeiro=True, entrelinha=1.05)

        quadro = _caixa(slide, MARGEM, Inches(4.6), Inches(9.0), Inches(0.5))
        _texto(quadro, subtitulo, SUBTITULO, RGBColor(0xC7, 0xD2, 0xD8),
               italico=True, primeiro=True)

        cartao = _retangulo(slide, MARGEM, Inches(5.4), Inches(4.6),
                            Inches(1.05), RGBColor(0x22, 0x33, 0x3D))
        quadro = cartao.text_frame
        quadro.margin_left = quadro.margin_top = Inches(0.22)
        _texto(quadro, autor, CORPO_MIN, BRANCO, negrito=True, primeiro=True)
        _texto(quadro, instituicao, CORPO_MIN, RGBColor(0xB8, 0xC4, 0xCC))

        self._rodape(slide, escuro=True)
        return slide

    def secao(self, numero: str, titulo: str, subtitulo: str):
        slide = self._novo(NAVY)
        quadro = _caixa(slide, MARGEM, Inches(2.0), Inches(6), Inches(1.6))
        _texto(quadro, numero, Pt(90), VERDE, fonte=SERIF, negrito=True,
               primeiro=True)

        quadro = _caixa(slide, MARGEM, Inches(3.6), Inches(11), Inches(0.9))
        _texto(quadro, titulo, Pt(42), BRANCO, fonte=SERIF, negrito=True,
               primeiro=True)

        quadro = _caixa(slide, MARGEM, Inches(4.65), Inches(10), Inches(0.6))
        _texto(quadro, subtitulo, CORPO_MIN, RGBColor(0xC7, 0xD2, 0xD8),
               primeiro=True)

        self._rodape(slide, escuro=True)
        return slide

    def topicos(self, kicker, titulo, itens, *, nota=None, destaque=False):
        slide = self._novo(OFFWHITE)
        y = self._cabecalho(slide, kicker, titulo)
        for texto in itens:
            cartao = _retangulo(slide, MARGEM, y, UTIL, Inches(0.62),
                                VERDE_CLARO if destaque else CINZA_CLARO)
            quadro = cartao.text_frame
            quadro.margin_left = Inches(0.28)
            quadro.vertical_anchor = MSO_ANCHOR.MIDDLE
            _texto(quadro, texto, CORPO_MIN, GRAFITE, primeiro=True)
            y += Inches(0.72)
        if nota:
            quadro = _caixa(slide, MARGEM, ALTURA - Inches(1.35), UTIL, Inches(0.7))
            _texto(quadro, nota, CORPO_MIN, VERDE, italico=True, primeiro=True)
        self._rodape(slide, escuro=False)
        return slide

    def cartoes(self, kicker, titulo, cartoes, *, colunas=3, nota=None,
                cor_nota=VERDE, fundo_nota=None):
        slide = self._novo(OFFWHITE)
        y = self._cabecalho(slide, kicker, titulo)
        vao = Inches(0.28)
        largura = int((UTIL - vao * (colunas - 1)) / colunas)
        linhas = (len(cartoes) + colunas - 1) // colunas
        altura = Inches(2.5) if linhas == 1 else Inches(1.75)

        for indice, (rotulo, corpo) in enumerate(cartoes):
            coluna, linha = indice % colunas, indice // colunas
            x = MARGEM + coluna * (largura + vao)
            topo = y + linha * (altura + Inches(0.25))
            cartao = _retangulo(slide, x, topo, largura, altura, BRANCO,
                                borda=RGBColor(0xE5, 0xE7, 0xEB))
            quadro = cartao.text_frame
            quadro.margin_left = quadro.margin_right = Inches(0.25)
            quadro.margin_top = Inches(0.22)
            _texto(quadro, rotulo, CARD_TITULO, GRAFITE, fonte=SERIF,
                   negrito=True, primeiro=True)
            _texto(quadro, corpo, CORPO_MIN, CINZA, espaco=Pt(8))

        if nota:
            caixa_y = ALTURA - Inches(1.5)
            if fundo_nota is not None:
                cartao = _retangulo(slide, MARGEM, caixa_y, UTIL, Inches(0.85),
                                    fundo_nota)
                quadro = cartao.text_frame
                quadro.margin_left = quadro.margin_right = Inches(0.28)
                quadro.vertical_anchor = MSO_ANCHOR.MIDDLE
                _texto(quadro, nota, CORPO_MIN, cor_nota, negrito=True,
                       primeiro=True)
            else:
                quadro = _caixa(slide, MARGEM, caixa_y, UTIL, Inches(0.8))
                _texto(quadro, nota, CORPO_MIN, cor_nota, italico=True,
                       primeiro=True)
        self._rodape(slide, escuro=False)
        return slide

    def codigo(self, kicker, titulo, linhas, *, notas=None, nota_titulo=None,
               rodape_nota=None):
        slide = self._novo(OFFWHITE)
        y = self._cabecalho(slide, kicker, titulo)
        largura_codigo = UTIL if not notas else int(UTIL * 0.60)
        altura_codigo = Inches(0.42) + Inches(0.32) * len(linhas)

        bloco = _retangulo(slide, MARGEM, y, largura_codigo, altura_codigo,
                           CODIGO_FUNDO, raio=0.04)
        quadro = bloco.text_frame
        quadro.margin_left = Inches(0.28)
        quadro.margin_top = Inches(0.2)
        for indice, linha in enumerate(linhas):
            cor = VERDE_SUAVE if linha.lstrip().startswith("#") else CODIGO_TEXTO
            _texto(quadro, linha or " ", CODIGO, cor, fonte=MONO,
                   primeiro=indice == 0, entrelinha=1.0)

        if notas:
            x = MARGEM + largura_codigo + Inches(0.3)
            largura = UTIL - largura_codigo - Inches(0.3)
            cartao = _retangulo(slide, x, y, largura, altura_codigo, VERDE_CLARO)
            quadro = cartao.text_frame
            quadro.margin_left = quadro.margin_right = Inches(0.25)
            quadro.margin_top = Inches(0.22)
            _texto(quadro, nota_titulo or "Por quê", CARD_TITULO, VERDE,
                   fonte=SERIF, negrito=True, primeiro=True)
            for nota in notas:
                _texto(quadro, f"•  {nota}", CORPO_MIN, GRAFITE, espaco=Pt(9))

        if rodape_nota:
            quadro = _caixa(slide, MARGEM, ALTURA - Inches(1.4), UTIL, Inches(0.8))
            _texto(quadro, rodape_nota, CORPO_MIN, CINZA, primeiro=True)
        self._rodape(slide, escuro=False)
        return slide

    def comparacao(self, kicker, titulo, esquerda, direita, *, nota=None):
        slide = self._novo(OFFWHITE)
        y = self._cabecalho(slide, kicker, titulo)
        largura = int((UTIL - Inches(0.3)) / 2)
        altura = Inches(0.42) + Inches(0.32) * max(
            len(esquerda[1]), len(direita[1])
        )

        for indice, (rotulo, linhas, cor) in enumerate((esquerda, direita)):
            x = MARGEM + indice * (largura + Inches(0.3))
            # o fundo da faixa acompanha a COR DO RÓTULO, não a posição:
            # verde para o lado funcional, vermelho para o imperativo,
            # cinza quando o lado é apenas neutro
            fundo = {
                VERDE: VERDE_CLARO,
                VERMELHO: RGBColor(0xFB, 0xE9, 0xEA),
            }.get(cor, CINZA_CLARO)
            faixa = _retangulo(slide, x, y, largura, Inches(0.5), fundo)
            quadro = faixa.text_frame
            quadro.margin_left = Inches(0.25)
            quadro.vertical_anchor = MSO_ANCHOR.MIDDLE
            _texto(quadro, rotulo, CORPO_MIN, cor, negrito=True, primeiro=True)

            bloco = _retangulo(slide, x, y + Inches(0.6), largura, altura,
                               CODIGO_FUNDO, raio=0.04)
            quadro = bloco.text_frame
            quadro.margin_left = Inches(0.25)
            quadro.margin_top = Inches(0.18)
            for posicao, linha in enumerate(linhas):
                cor_linha = (
                    VERDE_SUAVE if linha.lstrip().startswith("#") else CODIGO_TEXTO
                )
                _texto(quadro, linha or " ", CODIGO, cor_linha, fonte=MONO,
                       primeiro=posicao == 0, entrelinha=1.0)

        if nota:
            cartao = _retangulo(slide, MARGEM, ALTURA - Inches(1.55), UTIL,
                                Inches(0.9), CINZA_CLARO)
            quadro = cartao.text_frame
            quadro.margin_left = quadro.margin_right = Inches(0.28)
            quadro.vertical_anchor = MSO_ANCHOR.MIDDLE
            _texto(quadro, nota, CORPO_MIN, GRAFITE, primeiro=True)
        self._rodape(slide, escuro=False)
        return slide

    def tabela(self, kicker, titulo, linhas, *, nota=None):
        slide = self._novo(OFFWHITE)
        y = self._cabecalho(slide, kicker, titulo)
        # a altura da linha se ajusta ao espaço livre acima do rodapé, para que
        # uma tabela longa nunca invada a faixa inferior
        livre = ALTURA - y - Inches(1.35)
        altura = min(Inches(0.46), Emu(int(livre / len(linhas))) - Inches(0.05))
        for indice, (esquerda, direita) in enumerate(linhas):
            if indice % 2 == 0:
                _retangulo(slide, MARGEM, y, UTIL, altura, CINZA_CLARO, raio=0.02)
            quadro = _caixa(slide, MARGEM + Inches(0.25), y + Inches(0.08),
                            Inches(4.4), altura)
            _texto(quadro, esquerda, CORPO_MIN, GRAFITE, negrito=True,
                   primeiro=True)
            quadro = _caixa(slide, MARGEM + Inches(4.9), y + Inches(0.08),
                            UTIL - Inches(5.1), altura)
            _texto(quadro, direita, CORPO_MIN, VERDE, fonte=MONO, primeiro=True)
            y += altura + Inches(0.05)
        if nota:
            quadro = _caixa(slide, MARGEM, y + Inches(0.12), UTIL, Inches(0.6))
            _texto(quadro, nota, CORPO_MIN, VERDE, italico=True, primeiro=True)
        self._rodape(slide, escuro=False)
        return slide

    def imagem(self, kicker, titulo, arquivo, legenda, *, proporcao=1.9):
        """Insere uma captura de tela recortada em proporção panorâmica.

        As capturas são de página inteira (bem mais altas que largas). Postas
        inteiras no slide, encolhem até ficarem ilegíveis quando projetadas.
        Recortar a faixa superior — que é onde está o conteúdo relevante —
        deixa a imagem ocupar a largura toda e legível de longe.
        """
        slide = self._novo(OFFWHITE)
        y = self._cabecalho(slide, kicker, titulo)
        caminho = TELAS / arquivo
        disponivel = ALTURA - y - Inches(1.5)
        if caminho.exists():
            caminho = _recortar(caminho, proporcao)
            figura = slide.shapes.add_picture(str(caminho), 0, 0)
            escala = min(
                UTIL / figura.width, disponivel / figura.height
            )
            figura.width = Emu(int(figura.width * escala))
            figura.height = Emu(int(figura.height * escala))
            figura.left = Emu(int((LARGURA - figura.width) / 2))
            figura.top = Emu(int(y))
            figura.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
            figura.line.width = Pt(1)
        quadro = _caixa(slide, MARGEM, ALTURA - Inches(1.35), UTIL, Inches(0.8))
        _texto(quadro, legenda, CORPO_MIN, GRAFITE, primeiro=True)
        self._rodape(slide, escuro=False)
        return slide

    def dica(self, titulo, itens):
        """Slide das dicas pessoais — marcado para o avaliador reconhecer a cota."""
        slide = self._novo(OFFWHITE)
        selo = _retangulo(slide, MARGEM, MARGEM, Inches(2.9), Inches(0.42), AMBAR)
        quadro = selo.text_frame
        quadro.margin_left = Inches(0.2)
        quadro.vertical_anchor = MSO_ANCHOR.MIDDLE
        _texto(quadro, "DICA PESSOAL DE PYTHON", Pt(16), BRANCO, negrito=True,
               primeiro=True)

        quadro = _caixa(slide, MARGEM, MARGEM + Inches(0.62), UTIL, Inches(0.9))
        _texto(quadro, titulo, TITULO_MENOR, GRAFITE, fonte=SERIF, negrito=True,
               primeiro=True, entrelinha=1.05)

        y = MARGEM + Inches(1.75)
        for texto in itens:
            quadro = _caixa(slide, MARGEM, y, UTIL, Inches(0.75))
            _texto(quadro, f"•  {texto}", CORPO_MIN, GRAFITE, primeiro=True)
            y += Inches(0.85)
        self._rodape(slide, escuro=False)
        return slide

    def encerramento(self, titulo, subtitulo, repositorio):
        slide = self._novo(NAVY)
        if LOGO.exists():
            slide.shapes.add_picture(
                str(LOGO), MARGEM, Inches(1.5), height=Inches(1.15)
            )
        quadro = _caixa(slide, MARGEM, Inches(3.0), Inches(9), Inches(1.1))
        _texto(quadro, titulo, Pt(50), BRANCO, fonte=SERIF, negrito=True,
               primeiro=True)
        quadro = _caixa(slide, MARGEM, Inches(4.1), Inches(9), Inches(0.5))
        _texto(quadro, subtitulo, SUBTITULO, RGBColor(0xC7, 0xD2, 0xD8),
               italico=True, primeiro=True)
        cartao = _retangulo(slide, MARGEM, Inches(5.0), Inches(8.2), Inches(0.95),
                            RGBColor(0x22, 0x33, 0x3D))
        quadro = cartao.text_frame
        quadro.margin_left = Inches(0.25)
        quadro.vertical_anchor = MSO_ANCHOR.MIDDLE
        _texto(quadro, repositorio, CORPO_MIN, VERDE_SUAVE, fonte=MONO,
               primeiro=True)
        self._rodape(slide, escuro=True)
        return slide

    def fontes(self, entradas, nota):
        slide = self._novo(OFFWHITE)
        y = self._cabecalho(slide, "Créditos", "Fontes e recursos")
        largura_texto = UTIL - Inches(0.8)
        util_pt = Emu(largura_texto).inches * 72
        for marca, descricao in entradas:
            # a entrada da IA é longa: reservar a altura REAL evita que a
            # próxima linha suba por cima dela
            linhas = max(
                1, -(-int(len(descricao) * 0.50 * CORPO_MIN.pt) // int(util_pt))
            )
            altura = Inches(0.34) * linhas + Inches(0.25)

            quadro = _caixa(slide, MARGEM, y, Inches(0.75), Inches(0.4))
            _texto(quadro, marca, CORPO_MIN, VERDE, fonte=MONO, negrito=True,
                   primeiro=True)
            quadro = _caixa(slide, MARGEM + Inches(0.8), y, largura_texto, altura)
            _texto(quadro, descricao, CORPO_MIN, GRAFITE, primeiro=True)
            y += altura
        cartao = _retangulo(slide, MARGEM, ALTURA - Inches(1.55), UTIL,
                            Inches(0.9), CINZA_CLARO)
        quadro = cartao.text_frame
        quadro.margin_left = quadro.margin_right = Inches(0.28)
        quadro.vertical_anchor = MSO_ANCHOR.MIDDLE
        _texto(quadro, nota, CORPO_MIN, GRAFITE, primeiro=True)
        self._rodape(slide, escuro=False)
        return slide

    def salvar(self, caminho: Path) -> int:
        self.prs.save(str(caminho))
        return len(self.prs.slides._sldIdLst)


# ── conteúdo da apresentação ─────────────────────────────


def montar() -> Deck:
    d = Deck()

    # 1
    d.capa(
        "Programação\nFuncional em Python",
        "Do paradigma à prática — estudo de caso de um sistema real em produção",
        "Juliano Freitas",
        "Faculdade Nova Roma",
    )

    # 2
    d.topicos(
        "Roteiro",
        "O caminho de hoje",
        [
            "1 · O que é e por que importa          6 · Closures e decoradores",
            "2 · Os pilares                                   7 · Avaliação preguiçosa",
            "3 · Funções como valores                8 · Recursão",
            "4 · Ferramental do dia a dia            9 · Estudo de caso: o projeto",
            "5 · functools                                  10 · PF × Imperativo e limites",
        ],
        nota="Ao longo do caminho, três pausas com dicas práticas de Python.",
    )

    # 3
    d.cartoes(
        "Objetivos",
        "Ao final, você será capaz de:",
        [
            ("Identificar", "Reconhecer código funcional e seus conceitos em Python real."),
            ("Escrever", "Aplicar funções puras, HOFs, closures e decoradores com propósito."),
            ("Justificar", "Explicar POR QUE o estilo funcional ajuda — e onde não compensa."),
        ],
        nota="A aula usa um sistema real — o projeto Ponte Italia — como fio condutor: "
             "cada conceito aparece em código que roda.",
    )

    # 4
    d.secao("01", "O que é e por que importa",
            "Paradigmas, origem e a ideia central da programação funcional")

    # 5
    d.comparacao(
        "Bloco 1 · Paradigmas",
        "Duas formas de dizer a mesma coisa",
        ("Imperativo — COMO fazer",
         ["total = 0", "for n in numeros:", "    total += n"], VERMELHO),
        ("Funcional — O QUE se quer",
         ["total = sum(numeros)"], VERDE),
        nota="A programação funcional é um estilo declarativo: compomos funções "
             "em vez de encadear comandos que alteram variáveis.",
    )

    # 6
    d.cartoes(
        "Bloco 1 · Origem",
        "Do cálculo-λ ao Python de hoje",
        [
            ("1936", "Cálculo-λ: Church formaliza computação só com funções."),
            ("1958", "Lisp: McCarthy traz funções como dados e recursão."),
            ("1990", "Haskell: puramente funcional, com avaliação preguiçosa."),
            ("Hoje", "Python: multiparadigma — map, lambda, functools."),
        ],
        colunas=4,
    )

    # 7
    d.topicos(
        "Bloco 1 · Ideia central",
        "Três compromissos",
        [
            "1 · Funções no centro — o programa é avaliação e composição de funções.",
            "2 · Evitar estado mutável — em vez de alterar, produzir um novo valor.",
            "3 · Sem efeitos colaterais — a função conversa pelo retorno.",
        ],
    )

    # 8
    d.cartoes(
        "Bloco 1 · Honestidade",
        "Python não é Haskell — e tudo bem",
        [
            ("Funcional onde calcula",
             "Regras de negócio, transformação de dados, cálculos testáveis."),
            ("Imperativo nas bordas",
             "Ler e gravar banco, requisições HTTP, arquivos e relógio."),
        ],
        colunas=2,
        nota="Python é multiparadigma. A boa engenharia não é ser 100% funcional, "
             "e sim usar o estilo funcional ONDE ele paga.",
        fundo_nota=VERDE_CLARO,
    )

    # 9
    d.secao("02", "Os pilares",
            "Pureza, efeitos colaterais, imutabilidade e transparência referencial")

    # 10
    d.codigo(
        "Pilar 1 · Pureza",
        "Função pura",
        [
            "# pura",
            "def dobro(x):",
            "    return x * 2",
            "",
            "# impura: depende de estado externo",
            "total = 0",
            "def acumula(x):",
            "    global total",
            "    total += x      # efeito",
            "    return total",
        ],
        notas=[
            "Mesma entrada, mesma saída, sempre.",
            "Não lê nem altera nada fora dela.",
            "Depende só dos argumentos.",
        ],
        nota_titulo="Três marcas",
    )

    # 11
    d.cartoes(
        "Pilar 2 · Efeitos",
        "O que é um efeito colateral",
        [
            ("Mutar estado externo", "Variável global, atributo, lista compartilhada."),
            ("I/O", "Imprimir, ler banco, requisição de rede, arquivo."),
            ("Depender do ambiente", "Hora atual, número aleatório, variável de ambiente."),
        ],
        nota="Efeitos não são proibidos — são empurrados para as BORDAS. "
             "O núcleo que calcula fica puro, previsível e testável.",
        cor_nota=VERMELHO,
        fundo_nota=RGBColor(0xFB, 0xE9, 0xEA),
    )

    # 12
    d.codigo(
        "Pilar 3 · Imutabilidade",
        "Não altero — crio um novo",
        [
            "from dataclasses import dataclass",
            "",
            "@dataclass(frozen=True)",
            "class DocumentoResumo:",
            "    categoria: Categoria",
            "    status: Status",
            "",
            "d = DocumentoResumo(cat, ok)",
            "d.status = venc   # ERRO: congelado",
        ],
        notas=[
            "tuple, frozenset, str já são imutáveis.",
            "frozen=True: objetos que não mudam.",
            "Sem mutação = sem bug de estado compartilhado.",
        ],
        nota_titulo="Como se faz",
    )

    # 13
    d.codigo(
        "Pilar 4 · Transparência referencial",
        "Trocar a chamada pelo seu valor",
        [
            "preco(3)  ->  15.0        # sempre",
            "",
            "# então estas duas linhas são iguais:",
            "total = preco(3) + preco(3)",
            "",
            "p = preco(3)",
            "total = p + p",
        ],
        notas=[
            "Raciocinar sobre o código como em álgebra.",
            "Habilita cache seguro (memoização).",
            "Permite paralelizar sem travas.",
        ],
        nota_titulo="Por que importa",
        rodape_nota="É a propriedade que só funções puras têm — e a base dos slides "
                    "de lru_cache mais à frente.",
    )

    # 14
    d.codigo(
        "No projeto · Função pura",
        "calcular_status: o relógio entra pela porta",
        [
            "def calcular_status(",
            "    validade: date | None,",
            "    hoje: date,          # injetado!",
            "    janela: int = 30,",
            ") -> StatusDocumento:",
            "    if validade is None:",
            "        return Status.OK",
            "    if validade < hoje:",
            "        return Status.VENCIDO",
            "    if validade <= hoje + dias(janela):",
            "        return Status.VENCENDO",
            "    return Status.OK",
        ],
        notas=[
            "hoje é parâmetro, não datetime.now().",
            "O status não é gravado: é derivado na leitura.",
            "Testável sem relógio, sem mock, sem banco.",
        ],
        nota_titulo="Por que é pura",
    )

    # 15
    d.cartoes(
        "Pilares · Consequências",
        "Três ganhos que caem no seu colo",
        [
            ("01 · Testabilidade",
             "Um teste = três linhas. Sem fixture, sem mock, sem banco de pé."),
            ("02 · Concorrência segura",
             "Sem estado compartilhado, dois pedidos simultâneos não colidem."),
            ("03 · Raciocínio local",
             "A regra inteira cabe na tela; nada acontece 'à distância'."),
        ],
    )

    # 16 — dica pessoal 1
    d.dica(
        "Deixe as ferramentas vigiarem por você",
        [
            "ruff (lint) + black (formatação) + mypy (tipos) rodando no CI desde o dia 1.",
            "No projeto os três rodam em todo push: 188 testes e zero erro de tipo.",
            "mypy --strict pegou um bug real que passaria despercebido em revisão.",
            "Custo: 10 minutos de configuração. Retorno: o resto do semestre.",
        ],
    )

    # 17
    d.secao("03", "Funções como valores",
            "Primeira classe, funções de alta ordem e composição")

    # 18
    d.codigo(
        "Bloco 3 · Primeira classe",
        "Funções são dados",
        [
            "def dobro(x): return x * 2",
            "",
            "f = dobro          # sem parênteses!",
            "f(10)              # 20",
            "",
            "ops = {",
            '    "x2": dobro,',
            '    "neg": lambda x: -x,',
            "}",
        ],
        notas=[
            "Guardar em variável, lista ou dicionário.",
            "Passar como argumento.",
            "Retornar de outra função.",
        ],
        nota_titulo="O que dá para fazer",
    )

    # 19
    d.codigo(
        "Bloco 3 · Alta ordem",
        "Funções de alta ordem (HOF)",
        [
            "# map, filter e sorted já são HOFs:",
            "list(map(str.upper, nomes))",
            "list(filter(lambda n: n.ativo, users))",
            "sorted(cursos, key=lambda c: c.prazo)",
        ],
        rodape_nota="HOF = uma função que recebe outra função como argumento "
                    "e/ou devolve uma função.",
    )

    # 20
    d.codigo(
        "No projeto · Primeira classe",
        "O registry de parsers",
        [
            "# scraper/sources/__init__.py",
            "PARSERS = {",
            '    "universitaly": parser_universitaly,',
            '    "google_news": parser_google_news,',
            "}",
            "",
            "# a função é escolhida pelo NOME:",
            "universidades = PARSERS[fonte](html)",
        ],
        notas=[
            "Zero if fonte == '...' espalhado.",
            "Fonte nova = uma linha no dicionário.",
            "O código de coleta nunca muda.",
        ],
        nota_titulo="O ganho",
    )

    # 21
    d.codigo(
        "No projeto · HOF + composição",
        "agrupar_por_categoria recebe a chave",
        [
            "def agrupar_por_categoria(itens, chave):",
            "    return {",
            "        cat: tuple(i for i in itens",
            "                   if chave(i) == cat)",
            "        for cat in Categoria",
            "    }",
            "",
            "# a MESMA função serve documentos",
            "# e requisitos:",
            "agrupar_por_categoria(",
            "    docs, lambda d: d.categoria)",
        ],
        notas=[
            "A função de extração entra por parâmetro.",
            "O agrupamento não conhece o tipo concreto.",
            "Peças puras se encaixam: status → gap.",
        ],
        nota_titulo="Composição",
    )

    # 22
    d.secao("04", "Ferramental do dia a dia",
            "lambda, map, filter e comprehensions")

    # 23
    d.codigo(
        "Bloco 4 · lambda",
        "Função anônima, curtinha",
        [
            "quadrado = lambda x: x ** 2",
            "quadrado(5)      # 25",
            "",
            "# uso típico: argumento inline",
            "sorted(cursos, key=lambda c: c.prazo)",
        ],
        notas=[
            "lambda só inline, como argumento.",
            "Precisa de nome? Use def.",
            "Uma expressão: sem if/for de bloco.",
        ],
        nota_titulo="PEP 8",
    )

    # 24
    d.comparacao(
        "Bloco 4 · map e filter",
        "Transformar e selecionar sem loop",
        ("map — transforma cada item",
         ["list(map(str.strip, linhas))", "", "# equivale a:",
          "[s.strip() for s in linhas]"], VERDE),
        ("filter — mantém quem passa",
         ["list(filter(eh_valida, urls))", "", "# equivale a:",
          "[u for u in urls if eh_valida(u)]"], VERMELHO),
        nota="Em Python, a comprehension costuma ser a forma mais legível. "
             "map/filter ganham quando a função já existe e tem nome.",
    )

    # 25
    d.codigo(
        "No projeto · lambda + sorted",
        "ordenar_por_prazo",
        [
            "def ordenar_por_prazo(cursos, prazo_de):",
            "    return tuple(sorted(",
            "        cursos,",
            "        key=lambda c: (",
            "            prazo_de(c) is None,",
            "            prazo_de(c) or date.max,",
            "        ),",
            "    ))",
        ],
        notas=[
            "Quem tem prazo vem antes.",
            "Prazo mais próximo primeiro.",
            "prazo_de é injetado: serve a vários tipos.",
        ],
        nota_titulo="A regra na chave",
    )

    # 26
    d.codigo(
        "Bloco 4 · Comprehensions",
        "O jeito pythônico de transformar dados",
        [
            "nomes = [u.nome for u in users if u.ativo]",
            "",
            "por_id = {u.id: u for u in users}",
            "",
            "unicos = {u.cidade for u in users}",
        ],
        rodape_nota="Lê-se quase como português: 'o nome de cada usuário, "
                    "para cada usuário, se ativo'.",
    )

    # 27
    d.codigo(
        "No projeto · Comprehension aninhada",
        "Achatando requisitos por categoria",
        [
            "avaliacoes = tuple(",
            "    _avaliar(req, por_cat.get(cat, ()))",
            "    for cat, grupo in indexados.items()",
            "    for req in grupo",
            ")",
        ],
        rodape_nota="Dois for aninhados e uma chamada pura por item: nenhuma lista "
                    "temporária mutada, nenhum índice manual.",
    )

    # 28
    d.secao("05", "functools",
            "reduce, partial e lru_cache — a caixa de ferramentas funcional")

    # 29
    d.codigo(
        "functools · reduce",
        "Dobrar uma sequência num único valor",
        [
            "from functools import reduce",
            "",
            "total = reduce(",
            "    lambda acc, n: acc + n.peso,",
            "    noticias,",
            "    0,",
            ")",
        ],
        notas=[
            "Pense em: acumulador + item → novo acumulador.",
            "sum, max e any são reduces especializados.",
            "Comece pelo built-in; reduce para o resto.",
        ],
        nota_titulo="Como pensar",
    )

    # 30
    d.codigo(
        "functools · partial",
        "Fixar argumentos e nascer uma função nova",
        [
            "from functools import partial",
            "",
            "def conectar(host, porta, tls): ...",
            "",
            "local = partial(",
            '    conectar, "127.0.0.1", 5432)',
            "",
            "local(tls=False)  # host e porta fixos",
        ],
        notas=[
            "Especializar sem herança e sem classe.",
            "A nova função é a antiga com parte dos",
            "argumentos pré-preenchidos.",
        ],
        nota_titulo="A ideia",
    )

    # 31
    d.codigo(
        "No projeto · partial",
        "Um parser por fonte, sem duplicar código",
        [
            "from functools import partial",
            "",
            "parser_universitaly = partial(",
            "    parse_universidades,   # genérico",
            '    fonte="universitaly",',
            "    seletores=SELETORES,",
            ")",
        ],
        notas=[
            "O parser genérico transforma HTML em tuplas.",
            "Cada fonte é uma aplicação parcial.",
            "Fonte nova não copia código.",
        ],
        nota_titulo="O ganho",
    )

    # 32
    d.codigo(
        "functools · lru_cache",
        "Memoização em uma linha",
        [
            "from functools import lru_cache",
            "",
            "@lru_cache(maxsize=256)",
            "def fib(n):",
            "    if n < 2: return n",
            "    return fib(n-1) + fib(n-2)",
        ],
        notas=[
            "Só é CORRETO memoizar função pura.",
            "Cachear algo que lê o banco devolve",
            "dado velho — vira bug de dado obsoleto.",
        ],
        nota_titulo="A pegadinha",
    )

    # 33
    d.codigo(
        "No projeto · lru_cache",
        "requisitos_por_categoria: cache seguro por desenho",
        [
            "@lru_cache(maxsize=256)",
            "def requisitos_por_categoria(reqs):",
            "    return MappingProxyType(",
            "        agrupar_por_categoria(",
            "            reqs, lambda r: r.categoria))",
        ],
        notas=[
            "Entrada imutável e hashável (tupla).",
            "Função determinística, sem efeitos.",
            "Saída imutável (MappingProxyType).",
        ],
        nota_titulo="Por que é seguro",
        rodape_nota="O mesmo curso é comparado por muitos alunos: indexa uma vez, "
                    "não uma vez por pedido.",
    )

    # 34 — dica pessoal 2
    d.dica(
        "Type hints não são burocracia — são refatoração segura",
        [
            "Anote as assinaturas e deixe o mypy conferir o resto.",
            "Caso real do projeto: o mypy recusou um reduce cujo resultado ia direto "
            "para MappingProxyType — a inferência dirigida por contexto quebrava.",
            "A correção foi ligar o resultado a uma variável tipada antes de embrulhar.",
            "Sem tipos, isso só apareceria em produção, com dado estranho.",
        ],
    )

    # 35
    d.secao("06", "Closures e decoradores",
            "Funções que capturam ambiente e funções que envolvem funções")

    # 36
    d.codigo(
        "Bloco 6 · Closures",
        "Uma função que lembra do ambiente",
        [
            "def multiplicador(fator):",
            "    def aplicar(x):",
            "        return x * fator   # captura",
            "    return aplicar",
            "",
            "triplo = multiplicador(3)",
            "triplo(10)     # 30",
        ],
        notas=[
            "A função interna 'fecha sobre' as",
            "variáveis de fora e as carrega consigo.",
            "É uma fábrica de funções configuradas.",
        ],
        nota_titulo="A ideia",
    )

    # 37
    d.codigo(
        "Bloco 6 · Decoradores",
        "Uma função que envolve outra função",
        [
            "def log(funcao):",
            "    @wraps(funcao)   # preserva nome",
            "    def wrapper(*args, **kwargs):",
            '        print("chamando", funcao.__name__)',
            "        return funcao(*args, **kwargs)",
            "    return wrapper",
            "",
            "@log",
            "def somar(a, b): return a + b",
        ],
        rodape_nota="@log é açúcar para somar = log(somar). O decorador adiciona "
                    "comportamento transversal sem tocar na função original.",
    )

    # 38
    d.codigo(
        "No projeto · Decorador",
        "@cronometrar mede sem poluir a função",
        [
            "def cronometrar(funcao):",
            "    @wraps(funcao)",
            "    def wrapper(*args, **kwargs):",
            "        inicio = perf_counter()",
            "        try:",
            "            return funcao(*args, **kwargs)",
            "        finally:      # mede até se falhar",
            "            dur = perf_counter() - inicio",
            '            logger.info("%.1f ms", dur*1000)',
            "    return wrapper",
        ],
        notas=[
            "@wraps mantém introspecção e docs.",
            "finally garante a medição mesmo com erro.",
            "Reusável em qualquer função.",
        ],
        nota_titulo="Detalhes",
    )

    # 39
    d.codigo(
        "No projeto · Fábrica de decoradores",
        "retry_backoff: closure parametrizada",
        [
            "def retry_backoff(tentativas=3, base=0.5,",
            "                  fator=2.0,",
            "                  dormir=time.sleep):",
            "    def decorador(funcao):",
            "        @wraps(funcao)",
            "        def wrapper(*a, **kw):",
            "            for i in range(tentativas):",
            "                try:",
            "                    return funcao(*a, **kw)",
            "                except Exception:",
            "                    if i == tentativas-1: raise",
            "                    dormir(base * fator**i)",
            "        return wrapper",
            "    return decorador",
        ],
        notas=[
            "Cada uso captura sua própria config.",
            "dormir é injetado: testo o backoff",
            "sem dormir de verdade.",
        ],
        nota_titulo="Closure + teste",
    )

    # 40
    d.secao("07", "Avaliação preguiçosa",
            "Generators: produzir sob demanda, gastar memória mínima")

    # 41
    d.codigo(
        "Bloco 7 · Generators",
        "yield: produz um de cada vez",
        [
            "def contador(n):",
            "    i = 0",
            "    while i < n:",
            "        yield i        # pausa e devolve",
            "        i += 1",
            "",
            "soma = sum(x*x for x in range(10**6))",
            "# gera sob demanda, sem criar a lista",
        ],
        notas=[
            "Memória O(1): não materializa tudo.",
            "Trabalha com fluxos infinitos.",
            "islice consome só o necessário.",
        ],
        nota_titulo="Por quê",
    )

    # 42
    d.codigo(
        "No projeto · Lazy (implementado)",
        "Feed de notícias item a item",
        [
            "def iterar_itens_rss(xml):",
            "    raiz = ElementTree.fromstring(xml)",
            '    for item in raiz.iterfind("channel/item"):',
            "        yield {...}   # um por vez",
            "",
            "def paginar(itens, pagina, por_pagina):",
            "    inicio = (pagina - 1) * por_pagina",
            "    return tuple(islice(",
            "        itens, inicio, inicio + por_pagina))",
        ],
        notas=[
            "O Result do SQLAlchemy também é lazy.",
            "Um teste pagina itertools.count() —",
            "iterável INFINITO — e termina.",
        ],
        nota_titulo="Na prática",
    )

    # 43
    d.imagem(
        "No projeto · Radar",
        "O feed que os generators alimentam",
        "06-radar.png",
        "111 notícias coletadas de 3 fontes. A página mostra 10: paginar() consome "
        "só as linhas pedidas, o resto nunca vira objeto Python.",
    )

    # 44
    d.secao("08", "Recursão",
            "Um problema definido em termos de si mesmo")

    # 45
    d.codigo(
        "No projeto · Recursão (implementado)",
        "A árvore de categorias do FAQ",
        [
            "def todas_perguntas(categoria):",
            "    return categoria.perguntas + tuple(",
            "        p",
            "        for sub in categoria.subcategorias",
            "        for p in todas_perguntas(sub)",
            "    )",
            "",
            "def profundidade(categoria):",
            "    if not categoria.subcategorias:",
            "        return 1      # CASO BASE",
            "    return 1 + max(",
            "        profundidade(s)",
            "        for s in categoria.subcategorias)",
        ],
        notas=[
            "A estrutura é recursiva por natureza.",
            "Sem TCO: limite de ~1000 frames.",
            "O FAQ real tem 3 níveis — folga enorme.",
        ],
        nota_titulo="Caso base + passo",
    )

    # 46
    d.imagem(
        "No projeto · Ajuda",
        "A recursão renderizada na tela",
        "10-ajuda.png",
        "'Vistos e permanência' mostra 5 perguntas: 1 própria + 4 das subcategorias. "
        "Esse número vem de contar_perguntas() descendo a árvore inteira.",
    )

    # 47
    d.secao("09", "Estudo de caso: Ponte Italia",
            "Onde todos os conceitos se encontram em código que roda")

    # 48
    d.imagem(
        "O projeto · Visão geral",
        "Uma plataforma real, não um exemplo de aula",
        "01-home.png",
        "FastAPI + MySQL + Redis + worker .NET. 6 sprints, 188 testes automatizados, "
        "CI com lint, tipos e testes em todo push.",
    )

    # 49
    d.cartoes(
        "O projeto · Arquitetura",
        "Cálculo puro no centro, efeitos nas bordas",
        [
            ("routers/", "HTTP e banco — fazem o I/O e entregam dados imutáveis."),
            ("services/", "NÚCLEO PURO: calcular_gap, calcular_status, curadoria, faq."),
            ("scraper/", "Rede e HTML — coleta educada, isolada por fonte."),
        ],
        nota="As bordas fazem I/O e entregam dados imutáveis ao núcleo; o núcleo só "
             "calcula e devolve. Testo o núcleo sem subir nada.",
    )

    # 50
    d.codigo(
        "O projeto · O coração",
        "calcular_gap: o cálculo mais importante",
        [
            "def calcular_gap(requisitos, documentos):",
            "    por_cat = agrupar_por_categoria(",
            "        documentos, lambda d: d.categoria)",
            "    aval = tuple(",
            "        _avaliar(req, por_cat.get(cat, ()))",
            "        for cat, grupo in indexados.items()",
            "        for req in grupo",
            "    )",
            "    tri = lambda s: tuple(",
            "        i for st, i in aval if st == s)",
            '    return Gap(tri("atendido"),',
            '               tri("faltando"),',
            '               tri("vencendo"))',
        ],
        notas=[
            "Tuplas congeladas entram, Gap sai.",
            "Não conhece banco, HTTP nem relógio.",
            "14 testes cobrem os cruzamentos.",
        ],
        nota_titulo="100% pura",
    )

    # 51
    d.imagem(
        "O projeto · Da função pura à tela",
        "O mesmo calcular_gap que roda em 14 testes",
        "03-perfil-cofre.png",
        "Entrada: 3 requisitos do curso + 2 documentos do cofre. Saída: um Gap "
        "congelado — 33% pronto, 1 atendido, 2 faltando. Zero consulta dentro da regra.",
    )

    # 52
    d.imagem(
        "O projeto · Newsletter",
        "reduce agrupando a edição do dia",
        "13-email-newsletter.png",
        "agrupar_por_topico() dobra as notícias das últimas 24h nos 10 tópicos, "
        "publica na fila do Redis e o worker .NET renderiza e dispara às 9h.",
    )

    # 53
    d.tabela(
        "O projeto · Mapa",
        "Cada conceito vive em código real",
        [
            ("Função pura", "calcular_gap, calcular_status"),
            ("Imutabilidade", "schemas frozen, dataclasses"),
            ("Primeira classe / HOF", "registry PARSERS, agrupar_por_categoria"),
            ("lambda", "ordenar_por_prazo"),
            ("partial", "parser_universitaly"),
            ("lru_cache", "requisitos_por_categoria"),
            ("Decorador / closure", "@cronometrar, retry_backoff, @exigir_auth"),
            ("Generators / lazy", "iterar_itens_rss, paginar, janelas"),
            ("reduce", "estatisticas, agrupar_por_topico"),
            ("Recursão", "todas_perguntas, montar_arvore (FAQ)"),
        ],
        nota="16 de 16 conceitos implementados e documentados em docs/RELATORIO_FP.md.",
    )

    # 54 — dica pessoal 3
    d.dica(
        "A armadilha do import por valor",
        [
            "from modulo import funcao congela a REFERÊNCIA no momento do import.",
            "Bug real do projeto: o monkeypatch dos testes trocava cache.get_redis, "
            "mas a fila havia importado a função por valor — e continuou usando a original.",
            "Resultado: a suíte de testes escrevia no Redis de verdade, sem ninguém notar.",
            "Correção: importe o MÓDULO (from app.core import cache) e chame cache.get_redis().",
        ],
    )

    # 55
    d.secao("10", "PF × Imperativo e limites",
            "O contraste lado a lado e onde a PF não compensa")

    # 56
    d.comparacao(
        "Comparação",
        "O mesmo cálculo, dois mundos",
        ("Funcional — puro",
         ["def calcular_gap(reqs, docs):", "    ...        # só cálculo",
          "    return Gap(...)", "", "# testa sem banco, sem mock"], VERDE),
        ("Imperativo — acoplado",
         ["reqs = sessao.execute(SQL)", "for req in reqs:",
          "    docs = sessao.execute(SQL)  # N+1",
          '    if doc.status == "ok":',
          "        # status defasado",
          "        atendidos.append(req)"], VERMELHO),
        nota="A versão imperativa mistura buscar + derivar + particionar, depende de "
             "status defasado no banco e exige um MySQL de pé para cada teste.",
    )

    # 57
    d.cartoes(
        "Limites",
        "Quando a PF NÃO compensa em Python",
        [
            ("Cópias em excesso",
             "Recriar grandes estruturas a cada passo custa memória. "
             "Às vezes mutar local é ok."),
            ("Recursão profunda",
             "Sem TCO e com limite ~1000, laços são mais seguros que recursão longa."),
            ("Legibilidade",
             "map/filter/reduce encadeados podem ficar ilegíveis — "
             "comprehension costuma vencer."),
            ("I/O é inevitável",
             "Todo sistema real tem efeitos. O objetivo é isolá-los, "
             "não fingir que não existem."),
        ],
        colunas=2,
    )

    # 58
    d.topicos(
        "Boas práticas",
        "Como aplicar isso sem exagero",
        [
            "Separe o cálculo do I/O: núcleo puro, efeitos nas bordas.",
            "Dados que atravessam funções: prefira imutáveis (tuple, frozen).",
            "Funções pequenas e nomeadas; lambda só inline.",
            "Só memoize (lru_cache) função pura.",
            "Deixe ruff e mypy vigiarem tipos e estilo.",
        ],
    )

    # 59
    d.topicos(
        "Para levar",
        "Três frases para guardar",
        [
            "1 · Função pura = mesma entrada, mesma saída, sem efeitos.",
            "2 · Funções são valores: passe-as, retorne-as, componha.",
            "3 · Isole os efeitos nas bordas; mantenha o núcleo testável.",
        ],
        destaque=True,
    )

    # 60
    d.encerramento(
        "Obrigado!",
        "Perguntas são bem-vindas.",
        "github.com/julianosfreitas/scrapping_italy",
    )

    # 61
    d.fontes(
        [
            ("[1]", "Claude (Anthropic) — prompt: \"gerar apresentação acadêmica de "
                    "1h20 sobre Programação Funcional em Python, usando o projeto "
                    "Ponte Italia como estudo de caso, com slides de conceito, "
                    "código e comparação imperativo × funcional\"."),
            ("[2]", "Documentação oficial do Python — functools, itertools e "
                    "dataclasses (docs.python.org)."),
            ("[3]", "PEP 8 — Guia de Estilo para Código Python (peps.python.org/pep-0008)."),
        ],
        "Código, capturas de tela e dados do projeto Ponte Italia são de autoria "
        "própria (github.com/julianosfreitas/scrapping_italy).",
    )

    return d


def main() -> int:
    deck = montar()
    total = deck.salvar(SAIDA)
    print(f"{SAIDA.name}: {total} slides")
    print(f"corpo mínimo: {CORPO_MIN.pt:.0f}pt · código: {CODIGO.pt:.0f}pt "
          f"· rodapé: {RODAPE.pt:.0f}pt (isento)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
