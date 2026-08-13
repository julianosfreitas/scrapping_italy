"""Confere a geometria do .pptx: texto que vaza da caixa ou invade o rodapé.

Uso: ``python docs/verificar_layout.py [caminho.pptx]``

Complementa `verificar_slides.py`, que mede fontes e rodapé. Aqui o alvo é o
defeito visual que custa o ponto de "padrão que dificulta a leitura": linha de
código estourando o bloco escuro, cartão passando do rodapé, título largo
demais.

A medição é estimada a partir das métricas médias das fontes usadas — não
substitui olhar o arquivo, mas pega o que passa despercebido em 61 slides.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

from verificar_slides import eh_rodape

# largura média do glifo, em fração do corpo da fonte
AVANCO = {
    "Consolas": 0.550,  # monoespaçada: exato
    "Calibri": 0.500,
    "Georgia": 0.560,
}
AVANCO_PADRAO = 0.520

ALTURA_RODAPE = Pt(30)  # faixa inferior reservada ao rodapé
FOLGA = Pt(4)


def largura_estimada(texto: str, fonte: str | None, tamanho_pt: float) -> float:
    avanco = AVANCO.get(fonte or "", AVANCO_PADRAO)
    return len(texto) * avanco * tamanho_pt


def analisar(caminho: Path) -> list[str]:
    prs = Presentation(str(caminho))
    largura_slide = prs.slide_width
    altura_slide = prs.slide_height
    limite_inferior = altura_slide - ALTURA_RODAPE

    problemas: list[str] = []
    for numero, slide in enumerate(prs.slides, start=1):
        for forma in slide.shapes:
            if forma.left is None or forma.top is None:
                continue

            # forma decorativa (sem texto) pode sangrar de propósito
            if not forma.has_text_frame or not (forma.text_frame.text or "").strip():
                continue

            # o rodapé mora na faixa do rodapé — não é invasão
            if eh_rodape(forma.text_frame.text):
                continue

            direita = forma.left + (forma.width or 0)
            if direita > largura_slide + FOLGA or forma.left < -FOLGA:
                problemas.append(
                    f"slide {numero:2d}: caixa de texto sai pela lateral "
                    f"({Emu(direita).inches:.2f}in de {Emu(largura_slide).inches:.2f}in)"
                )

            quadro = forma.text_frame
            interno = (forma.width or 0) - quadro.margin_left - quadro.margin_right
            altura_texto = quadro.margin_top + quadro.margin_bottom

            for paragrafo in quadro.paragraphs:
                texto = "".join(c.text for c in paragrafo.runs)
                if not texto.strip():
                    altura_texto += Pt(12)
                    continue
                corrida = paragrafo.runs[0]
                tamanho = (corrida.font.size or Pt(18)).pt
                fonte = corrida.font.name
                espaco = paragrafo.line_spacing or 1.15
                altura_texto += Pt(tamanho * espaco)

                precisa = largura_estimada(texto, fonte, tamanho)
                # monoespaçada não quebra bem: qualquer estouro é visível
                if fonte == "Consolas" and Pt(precisa) > interno + FOLGA:
                    problemas.append(
                        f"slide {numero:2d}: código largo demais "
                        f"({Pt(precisa).inches:.2f}in > {Emu(interno).inches:.2f}in) "
                        f"-> {texto[:44]!r}"
                    )

            corpo_base = forma.top + altura_texto
            if corpo_base > limite_inferior + FOLGA and forma.has_text_frame:
                texto_inicial = (forma.text_frame.text or "").strip()[:40]
                if texto_inicial:
                    problemas.append(
                        f"slide {numero:2d}: texto invade o rodapé "
                        f"({Emu(corpo_base).inches:.2f}in > "
                        f"{Emu(limite_inferior).inches:.2f}in) -> {texto_inicial!r}"
                    )
    return problemas


def main(argumentos: list[str]) -> int:
    caminho = (
        Path(argumentos[0])
        if argumentos
        else Path(__file__).with_name("PF-Python-Ponte-Italia.pptx")
    )
    if not caminho.exists():
        print(f"arquivo não encontrado: {caminho}")
        return 2

    problemas = analisar(caminho)
    print(f"{caminho.name}: {len(problemas)} ocorrência(s) de layout\n")
    for linha in problemas:
        print(f"  {linha}")
    if problemas:
        print("\nResultado: revisar os slides acima.")
        return 1
    print("Resultado: nenhum estouro de caixa detectado.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
